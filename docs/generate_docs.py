"""
Generates the project deliverables for BistroSaaS:
    - er_diagram.png        (Entity-Relationship diagram, drawn with Pillow)
    - Project_Report.docx   (Word report with all required sections)
    - Project_Presentation.pptx (20-slide PowerPoint)

Run:  python docs/generate_docs.py
"""
import os

from PIL import Image, ImageDraw, ImageFont

DOCS_DIR = os.path.dirname(os.path.abspath(__file__))
ER_PATH = os.path.join(DOCS_DIR, "er_diagram.png")
DOCX_PATH = os.path.join(DOCS_DIR, "Project_Report.docx")
PPTX_PATH = os.path.join(DOCS_DIR, "Project_Presentation.pptx")


# ----------------------------------------------------------------------------
# Fonts helper
# ----------------------------------------------------------------------------
def load_font(size, bold=False):
    candidates = (
        ["arialbd.ttf", "Arial_Bold.ttf"] if bold else ["arial.ttf", "Arial.ttf"]
    )
    for name in candidates:
        for base in ["C:/Windows/Fonts/", "/usr/share/fonts/truetype/dejavu/", ""]:
            try:
                return ImageFont.truetype(base + name, size)
            except Exception:
                continue
    # DejaVu as a cross-platform fallback
    try:
        return ImageFont.truetype(
            "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf", size
        )
    except Exception:
        return ImageFont.load_default()


# ----------------------------------------------------------------------------
# 1) ER diagram (drawn with Pillow)
# ----------------------------------------------------------------------------
def build_er_diagram():
    W, H = 1240, 840
    img = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(img)

    title_font = load_font(28, bold=True)
    box_title_font = load_font(18, bold=True)
    attr_font = load_font(15)
    rel_font = load_font(15, bold=True)

    d.text((40, 20), "BistroSaaS - Entity Relationship Diagram",
           fill=(20, 20, 20), font=title_font)

    BLUE = (41, 128, 185)
    LIGHT = (236, 240, 241)
    DARK = (25, 25, 25)
    WIDTH = 300
    TITLE_H = 38
    LINE_H = 28

    entities = {
        "CustomerProfile": (40, 90,
            ["id (PK)", "user_id (FK, 1-1)", "phone", "address", "profile_pic"]),
        "User (Django built-in)": (470, 90,
            ["id (PK)", "username", "password", "email",
             "first_name", "last_name", "is_superuser"]),
        "Order": (900, 90,
            ["id (PK)", "user_id (FK)", "created_at", "total_amount",
             "status", "unseen_update", "updated_at"]),
        "MenuItem": (900, 530,
            ["id (PK)", "category_id (FK)", "name", "description",
             "price", "image", "is_available"]),
        "Category": (470, 600,
            ["id (PK)", "name", "description"]),
    }

    boxes = {}

    def draw_entity(name, x, y, attrs):
        height = TITLE_H + len(attrs) * LINE_H + 10
        d.rectangle([x, y, x + WIDTH, y + TITLE_H], fill=BLUE, outline=DARK, width=2)
        d.text((x + 12, y + 9), name, fill="white", font=box_title_font)
        d.rectangle([x, y + TITLE_H, x + WIDTH, y + height], fill=LIGHT,
                    outline=DARK, width=2)
        for i, attr in enumerate(attrs):
            d.text((x + 14, y + TITLE_H + 8 + i * LINE_H), attr,
                   fill=(30, 30, 30), font=attr_font)
        boxes[name] = (x, y, x + WIDTH, y + height)

    for name, (x, y, attrs) in entities.items():
        draw_entity(name, x, y, attrs)

    def center_y(box):
        return (box[1] + box[3]) // 2

    def rel_line(p1, p2, label):
        d.line([p1, p2], fill=(192, 57, 43), width=3)
        mx, my = (p1[0] + p2[0]) // 2, (p1[1] + p2[1]) // 2
        tw = d.textlength(label, font=rel_font)
        d.rectangle([mx - tw / 2 - 6, my - 14, mx + tw / 2 + 6, my + 14],
                    fill="white", outline=(192, 57, 43), width=2)
        d.text((mx - tw / 2, my - 9), label, fill=(192, 57, 43), font=rel_font)

    cp, usr, order = boxes["CustomerProfile"], boxes["User (Django built-in)"], boxes["Order"]
    item, cat = boxes["MenuItem"], boxes["Category"]

    # User 1-1 CustomerProfile
    rel_line((cp[2], center_y(cp)), (usr[0], center_y(usr)), "1 : 1")
    # User 1-M Order
    rel_line((usr[2], center_y(usr)), (order[0], center_y(order)), "1 : M")
    # Order M-N MenuItem (through order_items join table)
    rel_line(((order[0] + order[2]) // 2, order[3]),
             ((item[0] + item[2]) // 2, item[1]), "M : N")
    # Category 1-M MenuItem
    rel_line((item[0], center_y(item)), (cat[2], center_y(cat)), "M : 1")

    img.save(ER_PATH)
    print("Saved:", ER_PATH)


# ----------------------------------------------------------------------------
# 2) Word report
# ----------------------------------------------------------------------------
def build_report():
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor

    doc = Document()

    # ---- Title page ----
    for _ in range(3):
        doc.add_paragraph()
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = t.add_run("BistroSaaS")
    run.bold = True
    run.font.size = Pt(40)
    run.font.color.rgb = RGBColor(0x29, 0x80, 0xB9)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run("Restaurant Management & Online Ordering System")
    r.font.size = Pt(18)

    rep = doc.add_paragraph()
    rep.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rr = rep.add_run("Project Report")
    rr.italic = True
    rr.font.size = Pt(16)

    for _ in range(6):
        doc.add_paragraph()

    info = doc.add_paragraph()
    info.alignment = WD_ALIGN_PARAGRAPH.CENTER
    info.add_run(
        "Course: _______________________\n"
        "Instructor: _______________________\n"
        "Department: _______________________\n"
        "Submission Date: _______________________"
    ).font.size = Pt(13)

    doc.add_page_break()

    # ---- helper ----
    def heading(text):
        h = doc.add_heading(text, level=1)
        return h

    def bullets(items):
        for it in items:
            doc.add_paragraph(it, style="List Bullet")

    # ---- Group member information ----
    heading("Group Member Information")
    table = doc.add_table(rows=1, cols=4)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Name"
    hdr[1].text = "Roll / Reg. No."
    hdr[2].text = "Email"
    hdr[3].text = "Responsibilities"
    sample = [
        ("Mussa Afridi", "", "mussaafridi8240@gmail.com", "Backend, Deployment"),
        ("", "", "", ""),
        ("", "", "", ""),
        ("", "", "", ""),
    ]
    for name, roll, email, resp in sample:
        cells = table.add_row().cells
        cells[0].text = name
        cells[1].text = roll
        cells[2].text = email
        cells[3].text = resp
    doc.add_paragraph()
    doc.add_paragraph("(Fill in the remaining group members above.)").italic = True

    # ---- Introduction ----
    heading("Introduction")
    doc.add_paragraph(
        "BistroSaaS is a full-stack, web-based restaurant management and online "
        "ordering system developed using the Django web framework. It allows "
        "customers to browse a digital menu, search and filter dishes, place "
        "orders online, and track their order status in real time. At the same "
        "time, it provides restaurant administrators with powerful tools to "
        "manage menu items, categories, and incoming orders through a clean, "
        "responsive interface as well as the built-in Django admin panel."
    )
    doc.add_paragraph(
        "The application follows Django's Model-View-Template (MVT) architecture "
        "and includes a REST API (built with Django REST Framework) that exposes "
        "menu data in JSON format. The user interface is styled with Bootstrap 5 "
        "to ensure a modern, mobile-friendly experience."
    )

    # ---- Problem statement ----
    heading("Problem Statement")
    doc.add_paragraph(
        "Many small and medium restaurants still rely on manual, paper-based "
        "methods for taking and tracking orders. This leads to several problems:"
    )
    bullets([
        "Order errors caused by manual writing and verbal communication.",
        "No central, searchable record of the menu or past orders.",
        "Customers cannot view the menu or order online, limiting reach.",
        "Difficulty tracking the live status of an order (pending, preparing, completed).",
        "Menu updates (prices, availability, images) are slow and inconsistent.",
    ])
    doc.add_paragraph(
        "BistroSaaS solves these problems by digitizing the entire ordering "
        "workflow in a single, easy-to-use web application."
    )

    # ---- Objectives ----
    heading("Objectives")
    bullets([
        "Provide secure user registration and login with two roles: Customer and Admin.",
        "Allow customers to browse a categorized menu with search, filtering, and pagination.",
        "Enable customers to place orders online and view their order history.",
        "Allow admins to perform full CRUD operations on menu items (including image uploads).",
        "Implement an order confirmation workflow with status updates and customer notifications.",
        "Provide a user profile module to edit personal information and change passwords.",
        "Expose menu data through a RESTful API in JSON format.",
        "Deliver a responsive, modern UI using Bootstrap 5.",
    ])

    # ---- System design ----
    heading("System Design")
    doc.add_paragraph(
        "The system is built on Django's MVT (Model-View-Template) architecture:"
    )
    bullets([
        "Model (models.py): Defines the database tables - CustomerProfile, Category, MenuItem, and Order.",
        "View (views.py): Contains the logic for each page using both function-based and class-based views.",
        "Template (HTML files): Bootstrap 5 pages that present the data to the user.",
        "URLs (urls.py): Map web addresses to the correct views.",
        "Forms (forms.py): Validate and process user input.",
        "REST API (serializers.py + viewset): Serve menu data as JSON.",
    ])
    doc.add_paragraph("User Roles:", style="Intense Quote")
    bullets([
        "Guest: can browse the menu and read the API; must register/login to order.",
        "Customer: can place orders, view order history, and manage their profile.",
        "Admin (superuser): full access - menu management, order confirmation, and the Django admin panel.",
    ])
    doc.add_paragraph(
        "A key workflow is the order lifecycle: a customer places an order (status "
        "= Pending); an admin confirms and updates it (Preparing -> Completed, or "
        "Cancelled); the customer is then notified of the status change via an "
        "in-app notification badge."
    )

    # ---- Database design ----
    heading("Database Design (ER Diagram)")
    doc.add_paragraph(
        "The database is implemented in SQLite and consists of four main tables "
        "plus Django's built-in User table. Their relationships are:"
    )
    bullets([
        "User - CustomerProfile : One-to-One (each user has one profile).",
        "User - Order : One-to-Many (a user can place many orders).",
        "Category - MenuItem : One-to-Many (a category contains many items).",
        "Order - MenuItem : Many-to-Many (an order has many items; an item can be in many orders).",
    ])
    if os.path.exists(ER_PATH):
        doc.add_picture(ER_PATH, width=Inches(6.3))
        cap = doc.paragraphs[-1]
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        c = doc.add_paragraph("Figure 1: Entity-Relationship Diagram")
        c.alignment = WD_ALIGN_PARAGRAPH.CENTER
        c.runs[0].italic = True

    # ---- Screenshots ----
    heading("Screenshots")
    doc.add_paragraph(
        "Insert screenshots of the running application below. Recommended pages "
        "to capture (from your live site):"
    )
    shots = [
        "Home / Menu page (with search and category filters)",
        "Registration page (showing the Role dropdown and Secret Key)",
        "Login page",
        "Menu item detail page",
        "Add / Edit menu item form (admin)",
        "Place Order page (with live total)",
        "My Orders page (with status badges and notifications)",
        "Manage Orders page (admin order confirmation)",
        "User Profile page",
        "Django Admin panel",
        "REST API page (/api/menu/)",
    ]
    for i, s in enumerate(shots, 1):
        p = doc.add_paragraph()
        p.add_run(f"Screenshot {i}: {s}").bold = True
        ph = doc.add_paragraph("[ Paste screenshot here ]")
        ph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        ph.runs[0].italic = True

    # ---- Tools & technologies ----
    heading("Tools & Technologies Used")
    tools = [
        ("Python 3", "Core programming language"),
        ("Django 6.0.5", "Backend web framework (MVT architecture)"),
        ("Django REST Framework 3.17", "Builds the JSON REST API"),
        ("SQLite", "Database (file-based, db.sqlite3)"),
        ("Bootstrap 5", "Frontend CSS framework for responsive UI (via CDN)"),
        ("HTML5 + Django Templates", "Page structure and dynamic content"),
        ("Pillow", "Image handling for uploads (menu/profile pictures)"),
        ("Git & GitHub", "Version control and code hosting"),
        ("PythonAnywhere", "Cloud hosting / live deployment"),
        ("VS Code / Cursor", "Code editor / IDE"),
    ]
    ttable = doc.add_table(rows=1, cols=2)
    ttable.style = "Light Grid Accent 1"
    ttable.rows[0].cells[0].text = "Tool / Technology"
    ttable.rows[0].cells[1].text = "Purpose"
    for tool, purpose in tools:
        cells = ttable.add_row().cells
        cells[0].text = tool
        cells[1].text = purpose

    # ---- Conclusion ----
    heading("Conclusion")
    doc.add_paragraph(
        "BistroSaaS successfully demonstrates a complete, real-world restaurant "
        "management system built with Django. It covers the full software stack: "
        "database modeling, authentication and role-based access control, CRUD "
        "operations, file uploads, an order management workflow with notifications, "
        "a REST API, and a responsive Bootstrap interface. The project was version "
        "controlled with Git/GitHub and deployed live on PythonAnywhere."
    )
    doc.add_paragraph(
        "Through this project we gained practical experience in full-stack web "
        "development, the MVT design pattern, relational database design, and "
        "cloud deployment. Future enhancements could include online payment "
        "integration, email notifications, and a customer ratings/reviews system."
    )

    # ---- References ----
    heading("References")
    refs = [
        "Django Documentation - https://docs.djangoproject.com/",
        "Django REST Framework - https://www.django-rest-framework.org/",
        "Bootstrap 5 Documentation - https://getbootstrap.com/",
        "Pillow (PIL) Documentation - https://pillow.readthedocs.io/",
        "SQLite - https://www.sqlite.org/",
        "PythonAnywhere Help - https://help.pythonanywhere.com/",
        "Project Repository - https://github.com/afridi-mussa/Resturant-Management-System",
    ]
    for r in refs:
        doc.add_paragraph(r, style="List Number")

    doc.save(DOCX_PATH)
    print("Saved:", DOCX_PATH)


# ----------------------------------------------------------------------------
# 3) PowerPoint presentation
# ----------------------------------------------------------------------------
def build_presentation():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt

    # --- Theme palette (matches the BistroSaaS website: dark navbar + warm accent) ---
    DARK = RGBColor(0x1E, 0x2A, 0x38)    # deep slate (navbar feel)
    ACCENT = RGBColor(0xFF, 0x6B, 0x35)  # warm orange (food/brand accent)
    LIGHT = RGBColor(0xF5, 0xF7, 0xFA)   # content background
    TEXT = RGBColor(0x24, 0x31, 0x40)    # body text
    MUTED = RGBColor(0x8A, 0x97, 0xA5)   # footer/secondary
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    SOFT = RGBColor(0xFD, 0xE8, 0xDF)    # light orange wash
    FONT = "Segoe UI"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    BLANK = prs.slide_layouts[6]

    def _no_line(shape):
        shape.line.fill.background()

    def _fill(shape, color):
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        _no_line(shape)

    def _bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _run(p, text, size, color, bold=False, italic=False):
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = FONT
        return r

    def _footer(slide, idx):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.0), Inches(12.25), Pt(1.4))
        _fill(line, RGBColor(0xDD, 0xE3, 0xE9))
        ft = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(9), Inches(0.35))
        _run(ft.text_frame.paragraphs[0],
             "BistroSaaS  |  Restaurant Management System", 10, MUTED)
        nt = slide.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.2), Inches(0.35))
        np = nt.text_frame.paragraphs[0]
        np.alignment = PP_ALIGN.RIGHT
        _run(np, str(idx), 10, MUTED, bold=True)

    def _title_block(slide, title, emoji=""):
        # left accent stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH)
        _fill(stripe, ACCENT)
        # decorative soft circle top-right with emoji
        if emoji:
            circ = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(11.55), Inches(0.45), Inches(1.25), Inches(1.25))
            _fill(circ, SOFT)
            et = slide.shapes.add_textbox(Inches(11.55), Inches(0.5), Inches(1.25), Inches(1.15))
            ep = et.text_frame.paragraphs[0]
            ep.alignment = PP_ALIGN.CENTER
            _run(ep, emoji, 36, TEXT)
        # title
        tt = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(10.6), Inches(0.85))
        _run(tt.text_frame.paragraphs[0], title, 32, DARK, bold=True)
        # underline accent
        ul = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.28), Inches(2.5), Inches(0.09))
        _fill(ul, ACCENT)

    def add_content_slide(idx, title, bullets, emoji="", image=None):
        slide = prs.slides.add_slide(BLANK)
        _bg(slide, LIGHT)
        _title_block(slide, title, emoji)
        body_w = Inches(6.7) if image else Inches(11.8)
        body = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), body_w, Inches(5.0))
        tf = body.text_frame
        tf.word_wrap = True
        first = True
        for b in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(13)
            if b == "":
                _run(p, " ", 18, TEXT)
                continue
            _run(p, "\u25B8  ", 20, ACCENT, bold=True)   # ▸ marker
            _run(p, b, 20, TEXT)
        if image and os.path.exists(image):
            # white card behind the image
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.45), Inches(1.65),
                Inches(5.4), Inches(4.9))
            _fill(card, WHITE)
            slide.shapes.add_picture(image, Inches(7.6), Inches(2.6), width=Inches(5.1))
        _footer(slide, idx)
        return slide

    def add_banner_slide(title_lines, subtitle, tag):
        slide = prs.slides.add_slide(BLANK)
        _bg(slide, DARK)
        # accent side bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
        _fill(bar, ACCENT)
        # decorative ring
        ring = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(10.4), Inches(4.3), Inches(3.2), Inches(3.2))
        ring.fill.background()
        ring.line.color.rgb = ACCENT
        ring.line.width = Pt(2.5)
        # plate emoji
        et = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(3), Inches(1.2))
        _run(et.text_frame.paragraphs[0], "\U0001F37D", 60, WHITE)  # plate w/ cutlery
        # title
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.7), Inches(11), Inches(2.0))
        tf = tb.text_frame
        first = True
        for ln, sz in title_lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            _run(p, ln, sz, WHITE, bold=True)
        # accent line
        ul = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(4.55), Inches(3.0), Inches(0.1))
        _fill(ul, ACCENT)
        # subtitle
        st = slide.shapes.add_textbox(Inches(0.9), Inches(4.8), Inches(10.5), Inches(0.8))
        _run(st.text_frame.paragraphs[0], subtitle, 22, RGBColor(0xD7, 0xDE, 0xE5))
        # tag pill
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(5.7), Inches(3.5), Inches(0.55))
        _fill(pill, ACCENT)
        pp = pill.text_frame.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        pill.text_frame.word_wrap = True
        _run(pp, tag, 16, WHITE, bold=True)
        return slide

    # ---------- 1. Title ----------
    add_banner_slide(
        [("BistroSaaS", 54)],
        "Restaurant Management & Online Ordering System",
        "PROJECT PRESENTATION",
    )

    # ---------- 2. Group members ----------
    add_content_slide(2, "Group Members", [
        "Mussa Afridi  -  Backend & Deployment",
        "[ Member 2  -  role ]",
        "[ Member 3  -  role ]",
        "[ Member 4  -  role ]",
        "",
        "Course: ____________     Instructor: ____________",
    ], emoji="\U0001F465")

    # ---------- 3. Introduction ----------
    add_content_slide(3, "Introduction", [
        "A full-stack web application for restaurant ordering & management.",
        "Built with the Django web framework (Python).",
        "Customers browse a menu, place orders, and track status.",
        "Admins manage the menu and confirm orders.",
        "Modern, responsive UI using Bootstrap 5.",
    ], emoji="\U0001F4D6")

    # ---------- 4. Problem statement ----------
    add_content_slide(4, "Problem Statement", [
        "Manual, paper-based order taking causes errors.",
        "No central searchable menu or order records.",
        "Customers cannot view the menu or order online.",
        "Hard to track live order status.",
        "Menu / price / availability updates are slow.",
    ], emoji="\u2753")

    # ---------- 5. Objectives ----------
    add_content_slide(5, "Objectives", [
        "Secure registration & login (Customer / Admin roles).",
        "Browse menu with search, filter & pagination.",
        "Online order placement + order history.",
        "Admin CRUD on menu items with image uploads.",
        "Order confirmation workflow with notifications.",
        "REST API for menu data + responsive UI.",
    ], emoji="\U0001F3AF")

    # ---------- 6. Tools & technologies ----------
    add_content_slide(6, "Tools & Technologies", [
        "Python 3 + Django 6.0.5 (MVT framework).",
        "Django REST Framework (JSON API).",
        "SQLite database.",
        "Bootstrap 5 + HTML (frontend).",
        "Pillow (image uploads).",
        "Git & GitHub, deployed on PythonAnywhere.",
    ], emoji="\U0001F6E0")

    # ---------- 7. System architecture ----------
    add_content_slide(7, "System Architecture (MVT)", [
        "Model (models.py)  -  database tables.",
        "View (views.py)  -  page logic (FBV + CBV).",
        "Template (HTML)  -  Bootstrap pages shown to user.",
        "URLs (urls.py)  -  map addresses to views.",
        "Flow:  URL  ->  View  ->  Model (data)  ->  Template  ->  Browser.",
    ], emoji="\U0001F3D7")

    # ---------- 8. User roles ----------
    add_content_slide(8, "User Roles", [
        "Guest  -  browse menu & read API only.",
        "Customer  -  place orders, view history, edit profile.",
        "Admin (superuser)  -  full menu & order management + Django admin.",
        "Admin sign-up is protected by a secret key.",
    ], emoji="\U0001F511")

    # ---------- 9. Database design (ER diagram) ----------
    add_content_slide(9, "Database Design", [
        "4 main tables + Django User.",
        "User  1-1  CustomerProfile.",
        "User  1-M  Order.",
        "Category  1-M  MenuItem.",
        "Order  M-N  MenuItem.",
    ], emoji="\U0001F5C4", image=ER_PATH)

    # ---------- 10. Features overview ----------
    add_content_slide(10, "Features Overview", [
        "Authentication with two roles.",
        "Menu browsing: search, filter, pagination.",
        "Menu CRUD with image uploads (admin).",
        "Online ordering with live total.",
        "Order confirmation & notifications.",
        "User profiles + REST API.",
    ], emoji="\u2B50")

    # ---------- 11. Authentication & roles ----------
    add_content_slide(11, "Authentication & Roles", [
        "Register, Login, Logout (Django auth).",
        "Role dropdown: Customer (default) or Admin.",
        "Admin requires Secret Key (admin123).",
        "Profile auto-created for every user (signal).",
    ], emoji="\U0001F510")

    # ---------- 12. Menu management ----------
    add_content_slide(12, "Menu Management (CRUD)", [
        "Admins add / edit / delete menu items.",
        "Image upload for each dish (Pillow).",
        "Items grouped by category.",
        "Class-Based Views + access control (admin only).",
    ], emoji="\U0001F354")

    # ---------- 13. Search/filter/pagination ----------
    add_content_slide(13, "Search, Filter & Pagination", [
        "Search dishes by name or description (?q=).",
        "Filter by category buttons (?category=).",
        "6 items per page with Previous / Next controls.",
        "Implemented in MenuListView.",
    ], emoji="\U0001F50D")

    # ---------- 14. Order placement ----------
    add_content_slide(14, "Order Placement Flow", [
        "Customer selects items with checkboxes.",
        "Live JavaScript total updates instantly.",
        "Server recomputes the real total on submit.",
        "Order saved with status = Pending.",
    ], emoji="\U0001F6D2")

    # ---------- 15. Order confirmation & notifications ----------
    add_content_slide(15, "Order Confirmation & Notifications", [
        "Admin updates status: Pending -> Preparing -> Completed / Cancelled.",
        "Customer gets a red notification badge.",
        "Updated orders highlighted on 'My Orders'.",
        "Badge clears when the customer views orders.",
    ], emoji="\U0001F514")

    # ---------- 16. User profile ----------
    add_content_slide(16, "User Profile", [
        "Edit first name, last name, email.",
        "Edit phone, address, profile picture.",
        "Change password (new + confirm).",
        "Accessible from the navbar profile dropdown.",
    ], emoji="\U0001F464")

    # ---------- 17. REST API ----------
    add_content_slide(17, "REST API (Bonus)", [
        "Built with Django REST Framework.",
        "GET /api/menu/  -  list all items as JSON.",
        "GET /api/menu/<id>/  -  single item.",
        "Browsable API in the browser.",
        "Permission: read for all, write for logged-in users.",
    ], emoji="\U0001F517")

    # ---------- 18. Admin panel ----------
    add_content_slide(18, "Django Admin Panel", [
        "Manage Users, Categories, Menu Items, Orders.",
        "Search, filters, and inline editing.",
        "Custom red 'Cancel' button on all forms.",
        "Admins cannot see / delete their own account.",
    ], emoji="\u2699")

    # ---------- 19. Deployment ----------
    add_content_slide(19, "Version Control & Deployment", [
        "Source code hosted on GitHub.",
        "Live deployment on PythonAnywhere.",
        "Steps: clone -> venv -> install -> migrate -> configure web app.",
        "Live URL: mussa.pythonanywhere.com",
    ], emoji="\U0001F680")

    # ---------- 20. Conclusion / Thank you ----------
    add_banner_slide(
        [("Conclusion", 40), ("Thank You!", 30)],
        "A complete full-stack Django app  -  live demonstration to follow.",
        "QUESTIONS  &  DEMO",
    )

    prs.save(PPTX_PATH)
    print("Saved:", PPTX_PATH)


if __name__ == "__main__":
    build_er_diagram()
    build_report()
    build_presentation()
    print("All documents generated in:", DOCS_DIR)
